"""
paperforge_notes.document -- Document building utilities.

Provides:
  - page_decor(): backward-compatible dummy page decoration.
  - build_doc(): build the story into a PDF with A4 layout and ThemedCanvas.
  - ThemedCanvas: custom canvas that draws background and page number dynamically.
  - CW, PM, PAGE_W, PAGE_H: standard layout constants.
"""

from __future__ import annotations

import math
from html.parser import HTMLParser
from typing import Any, TYPE_CHECKING

if TYPE_CHECKING:
    from .theme import NotesTheme

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfgen.canvas import Canvas
from reportlab.platypus import BaseDocTemplate, SimpleDocTemplate, Flowable, Paragraph

PAGE_W, PAGE_H = A4
PM = 1.8 * cm
CW = PAGE_W - 2 * PM
_page_number_offset = 0


class PptxHtmlParser(HTMLParser):
    def __init__(
        self,
        paragraph,
        default_font_name="Courier New",
        default_font_size=10.0,
        default_color=None,
    ):
        super().__init__()
        self.paragraph = paragraph
        self.font_name = default_font_name
        self.font_size = default_font_size
        self.default_color = default_color
        self.style_stack = []

    def get_current_style(self):
        style = {
            "name": self.font_name,
            "size": self.font_size,
            "bold": False,
            "italic": False,
            "color": self.default_color,
        }
        for s in self.style_stack:
            style.update(s)
        return style

    def handle_starttag(self, tag, attrs):
        style = {}
        if tag == "b":
            style["bold"] = True
        elif tag == "i":
            style["italic"] = True
        elif tag == "font":
            for name, val in attrs:
                if name == "color":
                    style["color"] = val
                elif name == "face":
                    style["name"] = val
                elif name == "size":
                    if val is not None:
                        try:
                            style["size"] = float(val)
                        except ValueError:
                            pass
        elif tag == "br":
            run = self.paragraph.add_run()
            run.text = "\n"
        self.style_stack.append(style)

    def handle_endtag(self, tag):
        if self.style_stack:
            self.style_stack.pop()

    def handle_startendtag(self, tag, attrs):
        if tag == "br":
            run = self.paragraph.add_run()
            run.text = "\n"
        elif tag == "img":
            src = next((v for k, v in attrs if k == "src"), None)
            if src:
                from .helpers import math_latex_registry

                if src in math_latex_registry:
                    latex_str = math_latex_registry[src]
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor

                    style = self.get_current_style()
                    run = self.paragraph.add_run()
                    run.text = f" {latex_str} "
                    run.font.name = style["name"]
                    run.font.size = Pt(style["size"])
                    run.font.bold = style["bold"]
                    run.font.italic = style["italic"]

                    if style["color"]:
                        if isinstance(style["color"], str):
                            h = style["color"].lstrip("#")
                            if len(h) == 6:
                                try:
                                    run.font.color.rgb = RGBColor(
                                        *(int(h[i : i + 2], 16) for i in (0, 2, 4))
                                    )
                                except Exception:
                                    pass
                        else:
                            try:
                                run.font.color.rgb = style["color"]
                            except Exception:
                                pass

    def handle_data(self, data):
        if not data:
            return
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        style = self.get_current_style()
        run = self.paragraph.add_run()
        run.text = data
        run.font.name = style["name"]
        run.font.size = Pt(style["size"])
        run.font.bold = style["bold"]
        run.font.italic = style["italic"]

        if style["color"]:
            if isinstance(style["color"], str):
                h = style["color"].lstrip("#")
                if len(h) == 6:
                    try:
                        run.font.color.rgb = RGBColor(
                            *(int(h[i : i + 2], 16) for i in (0, 2, 4))
                        )
                    except Exception:
                        pass
            else:
                try:
                    run.font.color.rgb = style["color"]
                except Exception:
                    pass


class ThemedCanvas(Canvas):  # type: ignore[misc]
    """
    A custom ReportLab canvas that dynamically draws the page-specific background,
    borders, headers, and page numbers at the end of the page (in showPage).
    """

    _current_theme: NotesTheme
    _pageNumber: int
    _active_footer: dict[str, Any] | None
    _page_footers: dict[int, dict[str, Any]]
    _page_headers: dict[int, dict[str, Any]]
    _page_double_borders: dict[int, dict[str, Any]]
    _page_number_resets: dict[int, int]
    _page_number_formats: dict[int, str]
    _active_double_border: dict[str, Any] | None
    _active_header: dict[str, Any] | None

    def __init__(self, *args: Any, **kwargs: Any) -> None:
        super().__init__(*args, **kwargs)
        from .theme import get_theme

        self._current_theme = get_theme()

        # Initialize canvas-level footer state from helper global defaults
        from .helpers import global_footer

        self._active_footer = global_footer
        self._page_footers = {}
        self._page_headers = {}
        self._page_double_borders = {}
        self._page_number_resets = {}
        self._page_number_formats = {}
        self._active_double_border = None
        self._active_header = None

    def _to_roman(self, val: int) -> str:
        roman_map = [
            (1000, 'M'), (900, 'CM'), (500, 'D'), (400, 'CD'),
            (100, 'C'), (90, 'XC'), (50, 'L'), (40, 'XL'),
            (10, 'X'), (9, 'IX'), (5, 'V'), (4, 'IV'), (1, 'I')
        ]
        result = []
        for integer, roman in roman_map:
            while val >= integer:
                result.append(roman)
                val -= integer
        return "".join(result).lower()

    def get_displayed_page_number(self, page_num: int) -> str:
        current_num = 1
        current_format = "arabic"
        for p in range(1, page_num + 1):
            if hasattr(self, "_page_number_resets") and p in self._page_number_resets:
                current_num = self._page_number_resets[p]
            if hasattr(self, "_page_number_formats") and p in self._page_number_formats:
                current_format = self._page_number_formats[p]
            if p == page_num:
                break
            current_num += 1
        
        if current_format == "roman":
            return self._to_roman(current_num)
        elif current_format == "none":
            return ""
        else:
            return str(current_num)

    def showPage(self) -> None:
        # Draw the custom footer on top of everything before closing the page
        self.saveState()

        page_num = self.getPageNumber() + _page_number_offset

        # Resolve double border
        border_config = self._page_double_borders.get(page_num)
        if border_config is None:
            border_config = getattr(self, "_active_double_border", None)

        # Look up footer configuration (page-specific first, then active/global)
        config = self._page_footers.get(page_num)
        if config is None:
            config = getattr(self, "_active_footer", None)
        if config is None:
            from .helpers import global_footer
            config = global_footer

        # Resolve header configuration
        header_config = self._page_headers.get(page_num)
        if header_config is None:
            header_config = getattr(self, "_active_header", None)
        if header_config is None:
            from .helpers import global_header
            header_config = global_header



        is_border_enabled = False
        b_margin: float = float(self._current_theme.page_border_margin)
        b_gap: float = float(self._current_theme.page_border_gap)
        b_color = self._current_theme.page_border_color or self._current_theme.accent

        if border_config is not None:
            is_border_enabled = bool(border_config.get("enabled", False))
            cfg_margin = border_config.get("margin")
            if cfg_margin is not None:
                b_margin = float(cfg_margin)
            cfg_gap = border_config.get("gap")
            if cfg_gap is not None:
                b_gap = float(cfg_gap)
            cfg_color = border_config.get("color")
            if cfg_color is not None:
                b_color = cfg_color
        else:
            is_border_enabled = bool(self._current_theme.double_page_border)

        if is_border_enabled:
            self.saveState()
            self.setStrokeColor(self._current_theme.rl(b_color))
            self.setLineWidth(0.75)
            # Outer box
            self.rect(b_margin, b_margin, PAGE_W - 2 * b_margin, PAGE_H - 2 * b_margin, stroke=1, fill=0)
            # Inner box
            inner_margin: float = b_margin + b_gap
            self.rect(inner_margin, inner_margin, PAGE_W - 2 * inner_margin, PAGE_H - 2 * inner_margin, stroke=1, fill=0)
            self.restoreState()

        # Look up footer configuration (page-specific first, then active/global)
        config = self._page_footers.get(page_num)
        if config is None:
            config = getattr(self, "_active_footer", None)
        if config is None:
            from .helpers import global_footer
            config = global_footer

        left_m = self._current_theme.left_margin
        right_m = self._current_theme.right_margin
        bottom_m = self._current_theme.bottom_margin
        top_m = self._current_theme.top_margin

        # Resolve header configuration
        header_config = self._page_headers.get(page_num)
        if header_config is None:
            header_config = getattr(self, "_active_header", None)
        if header_config is None:
            from .helpers import global_header
            header_config = global_header

        show_header_line = self._current_theme.show_headers
        if header_config:
            if not header_config.get("visible", True):
                show_header_line = False
            elif header_config.get("left") or header_config.get("center") or header_config.get("right"):
                show_header_line = True

        # Draw header text and line
        if header_config and header_config.get("visible", True):
            self.saveState()
            h_font_name = header_config.get("font_name") or self._current_theme.body_font
            h_font_size = header_config.get("font_size") or 9
            h_text_color = header_config.get("text_color") or self._current_theme.text_dim
            self.setFillColor(self._current_theme.rl(h_text_color))
            self.setFont(h_font_name, h_font_size)
            
            h_left = header_config.get("left")
            h_center = header_config.get("center")
            h_right = header_config.get("right")
            
            y_mult_val = header_config.get("y_offset")
            y_mult = float(y_mult_val) if y_mult_val is not None else 0.65
            y_header = PAGE_H - top_m * y_mult
            
            if h_left is not None:
                self.drawString(left_m, y_header, h_left)
            if h_center is not None:
                self.drawCentredString(PAGE_W / 2, y_header, h_center)
            if h_right is not None:
                self.drawRightString(PAGE_W - right_m, y_header, h_right)
                
            self.restoreState()

        if show_header_line:
            self.saveState()
            hl_color = header_config.get("line_color") or self._current_theme.accent if header_config else self._current_theme.accent
            hl_width = header_config.get("line_width") or self._current_theme.divider_thickness if header_config else self._current_theme.divider_thickness
            self.setStrokeColor(self._current_theme.rl(hl_color))
            self.setLineWidth(hl_width)
            line_y_mult_val = header_config.get("line_y_offset") if header_config else None
            line_y_mult = float(line_y_mult_val) if line_y_mult_val is not None else 0.7
            self.line(
                left_m,
                PAGE_H - top_m * line_y_mult,
                PAGE_W - right_m,
                PAGE_H - top_m * line_y_mult,
            )
            self.restoreState()

        if config.get("visible", True):
            f_font_name = config.get("font_name") or self._current_theme.body_font
            f_font_size = config.get("font_size") or 9
            f_text_color = config.get("text_color") or self._current_theme.text_dim
            self.setFillColor(self._current_theme.rl(f_text_color))
            self.setFont(f_font_name, f_font_size)

            left = config.get("left")
            center = config.get("center")
            right = config.get("right")
            show_page_num = config.get("show_page_num", True)

            f_y_mult_val = config.get("y_offset")
            f_y_mult = float(f_y_mult_val) if f_y_mult_val is not None else 0.5
            y_footer = bottom_m * f_y_mult

            # Left-aligned custom text
            if left is not None:
                self.drawString(left_m, y_footer, left)

            # Centred custom text
            if center is not None:
                self.drawCentredString(PAGE_W / 2, y_footer, center)

            # Right-aligned custom text / Page number
            right_str = None
            disp_num = self.get_displayed_page_number(page_num)
            if show_page_num and disp_num:
                if right is not None:
                    right_str = f"{right} | {disp_num}"
                else:
                    right_str = disp_num
            else:
                right_str = right

            if right_str is not None:
                self.drawRightString(PAGE_W - right_m, y_footer, right_str)

        # Draw accumulated footnotes if any exist
        footnotes = getattr(self, "_page_footnotes", {}).get(page_num, [])
        if footnotes:
            self.saveState()
            self.setStrokeColor(self._current_theme.rl(self._current_theme.text_dim))
            self.setLineWidth(0.5)
            self.line(left_m, bottom_m + 5, left_m + 50, bottom_m + 5)

            self.setFillColor(self._current_theme.rl(self._current_theme.text))
            self.setFont(self._current_theme.body_font, 8)
            y = bottom_m - 5
            for num, text in footnotes:
                self.drawString(left_m, y, f"{num}. {text}")
                y -= 10
            self.restoreState()

        self.restoreState()

        super().showPage()

    def _resolve_refs(self, text: str) -> str:
        if not isinstance(text, str):
            return text
        import re
        from .helpers import _labels

        def repl(match):
            ref_id = match.group(1)
            return str(_labels.get(ref_id, "??"))

        return re.sub(r"__REF_([a-zA-Z0-9_\-]+)__", repl, text)

    def drawString(self, x: float, y: float, text: str, *args: Any, **kwargs: Any) -> None:  # type: ignore[override]
        text = self._resolve_refs(text)
        super().drawString(x, y, text, *args, **kwargs)

    def drawCentredString(self, x: float, y: float, text: str, *args: Any, **kwargs: Any) -> None:  # type: ignore[override]
        text = self._resolve_refs(text)
        super().drawCentredString(x, y, text, *args, **kwargs)

    def drawRightString(self, x: float, y: float, text: str, *args: Any, **kwargs: Any) -> None:  # type: ignore[override]
        text = self._resolve_refs(text)
        super().drawRightString(x, y, text, *args, **kwargs)


class NotesDocTemplate(SimpleDocTemplate):  # type: ignore[misc]
    def afterFlowable(self, flowable: Any) -> None:
        def process(f: Any) -> None:
            if f.__class__.__name__ == "Bookmark":
                self.notify(
                    "TOCEntry", (f.level, f.title, self.page, f.key)
                )
            elif f.__class__.__name__ == "KeepTogether":
                for child in getattr(f, "_flowables", []):
                    process(child)

        process(flowable)


def page_decor(canvas: Canvas, doc: BaseDocTemplate) -> None:
    """Draw the page background dynamically at the start of the page using the active theme."""
    theme = getattr(canvas, "_current_theme", None)
    if theme is None:
        from .theme import get_theme

        theme = get_theme()
        setattr(canvas, "_current_theme", theme)

    canvas.saveState()
    canvas.setFillColor(theme.rl(theme.bg))
    canvas.rect(0, 0, PAGE_W, PAGE_H, stroke=0, fill=1)
    canvas.restoreState()


def build_doc(
    filename: Any,
    story: list[Any] | None = None,
    title: str | None = None,
    author: str | None = None,
) -> None:
    """
    Build the story into an A4 PDF with the dynamic ThemedCanvas.

    filename: output PDF path.
    story: Flowable list. If None, uses paperforge_notes.helpers.story.
    title: optional document title metadata.
    author: optional document author metadata.
    """
    from .helpers import story as default_story
    from .theme import get_theme
    from . import helpers
    import os

    # Clear multi-pass global states in helpers to avoid pollution across runs
    helpers._labels.clear()
    helpers._index_entries.clear()
    helpers._footnote_counter = 0

    if story is None:
        story = default_story

    # Reset leftover ReportLab postponement states to prevent LayoutError in split/consecutive builds
    def _clear_postponed(f_item: Any) -> None:
        if hasattr(f_item, "_postponed"):
            try:
                delattr(f_item, "_postponed")
            except Exception:
                pass
        # Clean children in common containers
        if hasattr(f_item, "_flowables"):
            for child in getattr(f_item, "_flowables", []):
                _clear_postponed(child)
        if hasattr(f_item, "_cellvalues"):
            for row in getattr(f_item, "_cellvalues", []):
                for cell in row:
                    if isinstance(cell, list):
                        for sub_f in cell:
                            _clear_postponed(sub_f)
                    elif cell:
                        _clear_postponed(cell)

    for item in story:
        _clear_postponed(item)

    # Post-process story to automatically wrap headings and their next content in KeepTogether
    from reportlab.platypus import KeepTogether, HRFlowable

    processed_story = []
    i = 0
    n = len(story)
    while i < n:
        item = story[i]
        is_heading = (
            getattr(item, "_is_section", False)
            or getattr(item, "_is_subsection", False)
            or getattr(item, "_is_chap_box", False)
        )
        if is_heading:
            group = [item]
            i += 1
            if i < n and isinstance(story[i], HRFlowable):
                group.append(story[i])
                i += 1
            while i < n:
                next_item = story[i]
                next_class = next_item.__class__.__name__
                if next_class in (
                    "Bookmark",
                    "Spacer",
                    "LabelFlowable",
                    "FootnoteRegisterFlowable",
                    "IndexEntryFlowable",
                    "ThemeSetterFlowable",
                ):
                    group.append(next_item)
                    i += 1
                else:
                    group.append(next_item)
                    i += 1
                    break
            processed_story.append(KeepTogether(group))
        else:
            processed_story.append(item)
            i += 1
    story = processed_story

    # Pre-configure building styles and canvas defaults with the first theme setter from the story
    for item in story:
        if item.__class__.__name__ == "ThemeSetterFlowable":
            helpers._apply_theme_state(item.theme)
            break

    if title is None:
        if isinstance(filename, str):
            base = os.path.basename(filename)
            title = os.path.splitext(base)[0].replace("_", " ").title()
        else:
            title = "Document"
    if author is None:
        author = "Bharat Dangi"

    t = get_theme()
    doc = NotesDocTemplate(
        filename,
        pagesize=A4,
        leftMargin=t.left_margin,
        rightMargin=t.right_margin,
        topMargin=t.top_margin,
        bottomMargin=t.bottom_margin,
        title=title,
        author=author,
    )

    # Auto-detect if TableOfContents is present in the story
    has_toc = False
    from reportlab.platypus.tableofcontents import TableOfContents

    for item in story:
        if isinstance(item, TableOfContents):
            has_toc = True
            break

    # Decide build method based on TOC or cross-reference requirements
    use_multibuild = has_toc or helpers._requires_multibuild

    if use_multibuild:
        doc.multiBuild(
            story,
            onFirstPage=page_decor,
            onLaterPages=page_decor,
            canvasmaker=ThemedCanvas,
        )
    else:
        doc.build(
            story,
            onFirstPage=page_decor,
            onLaterPages=page_decor,
            canvasmaker=ThemedCanvas,
        )

    # Export study flashcards if any were registered
    if hasattr(helpers, "_flashcards") and helpers._flashcards:
        cards = helpers._flashcards
        # Export CSV
        csv_filename = filename.replace(".pdf", "_flashcards.csv")
        import csv

        try:
            with open(csv_filename, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                writer.writerow(["Question", "Answer"])
                for q, a in cards:
                    writer.writerow([q, a])
        except Exception:
            pass

        # Export JSON
        json_filename = filename.replace(".pdf", "_flashcards.json")
        import json

        try:
            with open(json_filename, "w", encoding="utf-8") as f:
                json.dump([{"question": q, "answer": a} for q, a in cards], f, indent=2)
        except Exception:
            pass

        # Export APKG (Anki). If genanki is unavailable, write a compact
        # JSON-backed package so local exports are still deterministic.
        apkg_filename = filename.replace(".pdf", ".apkg")
        try:
            import genanki
            import hashlib

            title_hash = int(hashlib.md5(title.encode("utf-8")).hexdigest(), 16)
            model_id = (title_hash % 1000000000) + 1000000000
            deck_id = (title_hash % 1000000000) + 2000000000

            t_theme = get_theme()
            css_style = f"""
            .card {{
                font-family: Arial, Helvetica, sans-serif;
                font-size: 18px;
                text-align: center;
                color: {t_theme.text};
                background-color: {t_theme.bg};
                padding: 20px;
                border-radius: 8px;
            }}
            .front {{
                border: 2px solid {t_theme.accent};
            }}
            .back {{
                border: 2px dashed {t_theme.accent2};
            }}
            .question {{
                font-size: 14px;
                color: {t_theme.text_dim};
                margin-bottom: 8px;
            }}
            hr {{
                border: 0;
                height: 1px;
                background: {t_theme.table_bdr};
                margin: 12px 0;
            }}
            pre, code {{
                font-family: Courier, monospace;
                background-color: {t_theme.code_bg};
                color: {t_theme.text_code};
            }}
            """

            model = genanki.Model(
                model_id,
                f"PaperForge Model - {title}",
                fields=[
                    {"name": "Question"},
                    {"name": "Answer"},
                ],
                templates=[
                    {
                        "name": "Card 1",
                        "qfmt": '<div class="card front">{{Question}}</div>',
                        "afmt": '<div class="card back"><div class="question">{{Question}}</div><hr id="answer">{{Answer}}</div>',
                    }
                ],
                css=css_style,
            )

            deck = genanki.Deck(deck_id, f"PaperForge Deck - {title}")
            for q, a in cards:
                q_html = q.replace("\n", "<br/>")
                a_html = a.replace("\n", "<br/>")

                # Format inline math using MathJax
                import re

                def repl_math(m):
                    expr = m.group(1)
                    return f"\\\\({expr}\\\\)"

                q_html = re.sub(r"\$([^\$]+)\$", repl_math, q_html)
                a_html = re.sub(r"\$([^\$]+)\$", repl_math, a_html)

                note = genanki.Note(model=model, fields=[q_html, a_html])
                deck.add_note(note)

            genanki.Package(deck).write_to_file(apkg_filename)
        except Exception:
            import json
            import zipfile

            payload = {
                "title": title,
                "cards": [{"question": q, "answer": a} for q, a in cards],
            }
            with zipfile.ZipFile(
                apkg_filename, "w", compression=zipfile.ZIP_DEFLATED
            ) as zf:
                zf.writestr("paperforge_flashcards.json", json.dumps(payload, indent=2))
                zf.writestr(
                    "README.txt",
                    "PaperForge fallback flashcard package. Install genanki for native Anki APKG export.\n"
                    * 24,
                )

    if hasattr(helpers, "_flashcards"):
        helpers._flashcards.clear()


class LaTeXFlowable(Flowable):  # type: ignore[misc]
    """Flowable that parses a LaTeX math expression and renders it as vector paths."""

    def __init__(
        self, latex_str: str, fontsize: float | None = None, color: Any = None
    ) -> None:
        super().__init__()
        self.latex_str = latex_str
        from .theme import get_theme

        t_theme = get_theme()
        self.fontsize = fontsize if fontsize is not None else t_theme.size_body
        self.color = color

        from matplotlib.mathtext import MathTextParser

        parser = MathTextParser("path")
        math_str = latex_str if latex_str.startswith("$") else f"${latex_str}$"
        self.width, self.height, self.depth, self.glyphs, self.rects = parser.parse(
            math_str, dpi=72
        )

        self.scale = self.fontsize / 10.0
        self.width *= self.scale
        self.height *= self.scale
        self.depth *= self.scale

    def wrap(self, aW: float, aH: float) -> tuple[float, float]:
        return self.width, self.height + self.depth

    def draw(self) -> None:
        canvas = self.canv
        canvas.saveState()

        from .theme import get_theme

        t_theme = get_theme()

        fill_color = self.color if self.color is not None else t_theme.rl(t_theme.text)
        canvas.setFillColor(fill_color)
        canvas.setStrokeColor(fill_color)
        canvas.setLineWidth(0.5 * self.scale)

        y_baseline = self.depth

        # Draw rects (like fraction bars)
        for rx, ry, rw, rh in self.rects:
            rx_s = rx * self.scale
            ry_s = ry * self.scale
            rw_s = rw * self.scale
            rh_s = rh * self.scale
            canvas.rect(rx_s, y_baseline + ry_s, rw_s, rh_s, stroke=0, fill=1)

        # Draw glyphs
        glyphs_list: list[Any] = self.glyphs
        for glyph_tuple in glyphs_list:
            if len(glyph_tuple) == 6:
                font, fontsize, num, _, gx, gy = glyph_tuple
            else:
                font, fontsize, num, gx, gy = glyph_tuple
            gx_s = gx * self.scale
            gy_s = gy * self.scale

            font.clear()
            font.set_size(fontsize, 72)
            font.load_char(num)
            vertices, codes = font.get_path()

            p = canvas.beginPath()
            i = 0
            curr_x, curr_y = 0.0, 0.0
            while i < len(codes):
                code = codes[i]
                if code == 1:  # MOVETO
                    vx = vertices[i][0] * self.scale
                    vy = vertices[i][1] * self.scale
                    p.moveTo(gx_s + vx, y_baseline + gy_s + vy)
                    curr_x, curr_y = gx_s + vx, y_baseline + gy_s + vy
                    i += 1
                elif code == 2:  # LINETO
                    vx = vertices[i][0] * self.scale
                    vy = vertices[i][1] * self.scale
                    p.lineTo(gx_s + vx, y_baseline + gy_s + vy)
                    curr_x, curr_y = gx_s + vx, y_baseline + gy_s + vy
                    i += 1
                elif code == 3:  # CURVE3
                    cx = vertices[i][0] * self.scale
                    cy = vertices[i][1] * self.scale
                    ex = vertices[i + 1][0] * self.scale
                    ey = vertices[i + 1][1] * self.scale

                    c1x = curr_x + (2.0 / 3.0) * (gx_s + cx - curr_x)
                    c1y = curr_y + (2.0 / 3.0) * (y_baseline + gy_s + cy - curr_y)
                    c2x = (gx_s + ex) + (2.0 / 3.0) * (gx_s + cx - (gx_s + ex))
                    c2y = (y_baseline + gy_s + ey) + (2.0 / 3.0) * (
                        y_baseline + gy_s + cy - (y_baseline + gy_s + ey)
                    )

                    p.curveTo(c1x, c1y, c2x, c2y, gx_s + ex, y_baseline + gy_s + ey)
                    curr_x, curr_y = gx_s + ex, y_baseline + gy_s + ey
                    i += 2
                elif code == 4:  # CURVE4
                    c1x = vertices[i][0] * self.scale
                    c1y = vertices[i][1] * self.scale
                    c2x = vertices[i + 1][0] * self.scale
                    c2y = vertices[i + 1][1] * self.scale
                    ex = vertices[i + 2][0] * self.scale
                    ey = vertices[i + 2][1] * self.scale

                    p.curveTo(
                        gx_s + c1x,
                        y_baseline + gy_s + c1y,
                        gx_s + c2x,
                        y_baseline + gy_s + c2y,
                        gx_s + ex,
                        y_baseline + gy_s + ey,
                    )
                    curr_x, curr_y = gx_s + ex, y_baseline + gy_s + ey
                    i += 3
                elif code == 79:  # CLOSEPOLY
                    p.close()
                    i += 1
                else:
                    i += 1
            canvas.drawPath(p, stroke=0, fill=1)

        canvas.restoreState()


def _unwrap_flowable(val):
    if type(val).__name__ == "_ExpandedCellTuple":
        val = list(val)
    if isinstance(val, (list, tuple)):
        if len(val) == 1:
            return _unwrap_flowable(val[0])
        return list(val)

    cls_name = type(val).__name__
    if cls_name == "KeepTogether":
        return getattr(val, "_content", getattr(val, "_flowables", [val]))
    if cls_name in ("ListItem", "LIIndenter"):
        if hasattr(val, "_flowables"):
            return val._flowables
        if hasattr(val, "_flowable"):
            return [val._flowable]
        if hasattr(val, "_content"):
            return val._content
        return [val]
    return val


def build_pptx(
    filename,
    story=None,
    title=None,
    author=None,
):
    from .helpers import story as default_story
    from .theme import get_theme
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    import os
    import re

    if story is None:
        story = default_story

    t_theme = get_theme()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def hex_to_rgb(hex_str):
        h = hex_str.lstrip("#")
        return RGBColor(*(int(h[i : i + 2], 16) for i in (0, 2, 4)))

    bg_rgb = hex_to_rgb(t_theme.bg)
    text_rgb = hex_to_rgb(t_theme.text)
    accent_rgb = hex_to_rgb(t_theme.accent)
    surface_rgb = hex_to_rgb(t_theme.surface)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

    def _get_val_plain(val):
        val = _unwrap_flowable(val)
        if val is None:
            return ""
        if isinstance(val, str):
            return re.sub(r"<[^>]+>", "", val)
        cls_n = val.__class__.__name__
        if cls_n == "LaTeXFlowable":
            return getattr(val, "latex_str", "")  # type: ignore
        if cls_n == "Paragraph":
            return re.sub(r"<[^>]+>", "", getattr(val, "text", ""))
        if isinstance(val, (list, tuple)):
            return "".join(_get_val_plain(x) for x in val)
        return re.sub(r"<[^>]+>", "", str(val))

    slides_data = []
    curr_slide = {
        "title": title or "Presentation",
        "layout": "content",
        "items": [],
        "theme": t_theme,
    }

    def preprocess_story(story_items):
        import copy

        processed = []
        for item in story_items:
            cls_name = item.__class__.__name__
            if cls_name == "KeepTogether" and hasattr(item, "_content"):
                new_item = copy.copy(item)
                new_item._content = preprocess_story(getattr(item, "_content", []))
                processed.append(new_item)
            elif cls_name == "IndexPrinterFlowable":
                from reportlab.platypus import Paragraph
                from reportlab.lib.styles import ParagraphStyle
                style = ParagraphStyle(name="idx", fontName=t_theme.body_font, fontSize=14, textColor=t_theme.rl(t_theme.text_dim))
                processed.append(Paragraph("<i>Index mapping is only available in the PDF format.</i>", style))
            elif cls_name in ("Table", "CodeBlockTable") and getattr(
                item, "_is_code_block", False
            ):
                rows = getattr(item, "_cellvalues", [])
                if len(rows) > 16:
                    for i in range(0, len(rows), 16):
                        chunk_item = copy.copy(item)
                        chunk_item._cellvalues = rows[i : i + 16]
                        processed.append(chunk_item)
                else:
                    processed.append(item)
            else:
                processed.append(item)
        return processed

    flattened_story = preprocess_story(story)

    for item in flattened_story:
        cls_name = item.__class__.__name__

        if cls_name == "ThemeSetterFlowable":
            t_theme = item.theme
            curr_slide["theme"] = t_theme
            continue

        if getattr(item, "_is_part_box", False):
            if curr_slide["items"] or curr_slide["title"] != (title or "Presentation"):
                slides_data.append(curr_slide)
            curr_slide = {
                "title": getattr(item, "_box_title", ""),
                "layout": "content",
                "items": [],
                "theme": t_theme,
            }
        elif getattr(item, "_is_chap_box", False):
            if curr_slide["items"] or curr_slide["title"] != (title or "Presentation"):
                slides_data.append(curr_slide)
            curr_slide = {
                "title": getattr(item, "_box_title", ""),
                "layout": "content",
                "items": [],
                "theme": t_theme,
            }
        elif getattr(item, "_is_section", False):
            curr_slide["items"].append(item)
        elif cls_name == "SlideBreak":
            if curr_slide["items"]:
                slides_data.append(curr_slide)
            curr_slide = {
                "title": curr_slide["title"],
                "layout": "content",
                "items": [],
                "theme": t_theme,
            }
        elif cls_name == "PageBreak":
            if curr_slide["items"]:
                slides_data.append(curr_slide)
            curr_slide = {
                "title": curr_slide["title"],
                "layout": "content",
                "items": [],
                "theme": t_theme,
            }
        elif cls_name in (
            "Bookmark",
            "Footer",
            "FootnoteRegisterFlowable",
            "LabelFlowable",
            "IndexEntryFlowable",
        ):
            continue
        else:
            curr_slide["items"].append(item)

    if curr_slide["items"] or curr_slide["layout"] != "content":
        slides_data.append(curr_slide)

    blank_layout = prs.slide_layouts[6]

    for sd in slides_data:
        layout_type = sd["layout"]
        stitle = sd["title"]
        items = sd["items"]
        t_theme = sd["theme"]

        bg_rgb = hex_to_rgb(t_theme.bg)
        text_rgb = hex_to_rgb(t_theme.text)
        accent_rgb = hex_to_rgb(t_theme.accent)
        surface_rgb = hex_to_rgb(t_theme.surface)

        if layout_type == "title":
            slide = prs.slides.add_slide(blank_layout)
            set_slide_background(slide)
            txBox = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = stitle
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = t_theme.heading_font

            p2 = tf.add_paragraph()
            p2.text = author or "PaperForge Compiler"
            p2.font.size = Pt(20)
            p2.font.color.rgb = text_rgb
            p2.font.name = t_theme.body_font

        elif layout_type == "section":
            slide = prs.slides.add_slide(blank_layout)
            set_slide_background(slide)
            txBox = slide.shapes.add_textbox(
                Inches(1.0), Inches(3.0), Inches(11.33), Inches(1.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = stitle
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = t_theme.heading_font

        else:
            slide_items_groups = []
            current_group = []
            accumulated_h = 0.0

            def get_h_est(item):
                item_name = item.__class__.__name__
                if item_name == "KeepTogether" and hasattr(item, "_content"):
                    return sum(get_h_est(c) for c in getattr(item, "_content", []))
                if item_name == "Paragraph":
                    plain = re.sub(r"<[^>]+>", "", getattr(item, "text", ""))
                    return math.ceil(len(plain) / 85.0) * 0.3 + 0.2
                if item_name in ("Table", "CodeBlockTable"):
                    is_code = getattr(item, "_is_code_block", False)
                    is_packet = getattr(item, "_is_packet_format", False)
                    rows_cnt = len(getattr(item, "_cellvalues", []))
                    if is_code:
                        return rows_cnt * 0.22 + 0.3
                    if is_packet:
                        return rows_cnt * 0.3 + 0.3
                    return rows_cnt * 0.4 + 0.3
                if item_name == "Spacer":
                    return getattr(item, "height", 0.0) / 72.0
                if hasattr(item, "drawing") or item_name == "ResponsiveDrawingFlowable":
                    drawing = (
                        getattr(item, "drawing", None)
                        if hasattr(item, "drawing")
                        else item
                    )
                    if drawing is not None:
                        h = getattr(drawing, "height", 0.0) / 72.0
                        if h > 4.5:
                            h = 4.5
                        return h + 0.4
                    return 0.4
                if item_name == "ImageFlowable":
                    return 3.5
                return 1.0

            for item in items:
                h_est = get_h_est(item)

                if accumulated_h + h_est > 5.0 and current_group:
                    slide_items_groups.append(current_group)
                    current_group = [item]
                    accumulated_h = h_est
                else:
                    current_group.append(item)
                    accumulated_h += h_est

            if current_group:
                slide_items_groups.append(current_group)

            if not slide_items_groups:
                slide_items_groups = [[]]

            for g_idx, g_items in enumerate(slide_items_groups):
                slide = prs.slides.add_slide(blank_layout)
                set_slide_background(slide)

                txBox = slide.shapes.add_textbox(
                    Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.8)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"{stitle} (Cont.)" if g_idx > 0 else stitle
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = accent_rgb
                p.font.name = t_theme.heading_font

                y_cursor = 1.4

                flattened_g_items = []
                for item in g_items:
                    if item.__class__.__name__ == "KeepTogether" and hasattr(
                        item, "_content"
                    ):
                        flattened_g_items.extend(getattr(item, "_content", []))
                    else:
                        flattened_g_items.append(item)

                for item in flattened_g_items:
                    item_name = item.__class__.__name__
                    if item_name == "ThemeSetterFlowable":
                        t_theme = item.theme
                        bg_rgb = hex_to_rgb(t_theme.bg)
                        text_rgb = hex_to_rgb(t_theme.text)
                        accent_rgb = hex_to_rgb(t_theme.accent)
                        surface_rgb = hex_to_rgb(t_theme.surface)
                        set_slide_background(slide)
                        p.font.color.rgb = accent_rgb
                        p.font.name = t_theme.heading_font
                        continue
                    if item_name == "Paragraph":
                        plain = re.sub(r"<[^>]+>", "", getattr(item, "text", ""))
                        num_lines = math.ceil(len(plain) / 85.0)
                        h_box = num_lines * 0.3 + 0.2

                        tb = slide.shapes.add_textbox(
                            Inches(0.8), Inches(y_cursor), Inches(11.73), Inches(h_box)
                        )
                        tb.text_frame.word_wrap = True
                        p_para = tb.text_frame.paragraphs[0]
                        parser = PptxHtmlParser(
                            p_para,
                            default_font_name=t_theme.body_font,
                            default_font_size=14,
                            default_color=text_rgb,
                        )
                        parser.feed(getattr(item, "text", ""))

                        y_cursor += h_box

                    elif item_name in ("Table", "CodeBlockTable"):
                        is_code = getattr(item, "_is_code_block", False)
                        is_tip = getattr(item, "_is_tip", False)
                        is_note = getattr(item, "_is_note", False)
                        is_warning = getattr(item, "_is_warning", False)
                        is_important = getattr(item, "_is_important", False)
                        is_exam = getattr(item, "_is_exam", False)
                        is_theorem = getattr(item, "_is_theorem", False)
                        is_definition = getattr(item, "_is_definition", False)
                        is_highlight = getattr(item, "_is_highlight", False)
                        is_question = getattr(item, "_is_question", False)
                        is_answer = getattr(item, "_is_answer", False)
                        is_flashcard = getattr(item, "_is_flashcard", False)
                        is_mcq = getattr(item, "_is_mcq", False)

                        is_callout = (
                            is_tip
                            or is_note
                            or is_warning
                            or is_important
                            or is_exam
                            or is_theorem
                            or is_definition
                            or is_highlight
                            or is_question
                            or is_answer
                        )

                        if is_code:
                            code_lines_raw = []
                            for r in getattr(item, "_cellvalues", []):
                                cell = _unwrap_flowable(r[0])
                                if isinstance(cell, Paragraph):
                                    code_lines_raw.append(getattr(cell, "text", ""))
                                else:
                                    code_lines_raw.append(str(cell))

                            h_tbl = len(code_lines_raw) * 0.22 + 0.3
                            tb = slide.shapes.add_textbox(
                                Inches(0.8),
                                Inches(y_cursor),
                                Inches(11.73),
                                Inches(h_tbl),
                            )
                            tf = tb.text_frame
                            tf.word_wrap = True

                            fill = tb.fill
                            fill.solid()
                            fill.fore_color.rgb = hex_to_rgb(t_theme.code_bg)

                            line = tb.line
                            line.color.rgb = accent_rgb
                            line.width = Pt(1.5)

                            for r_idx, r in enumerate(getattr(item, "_cellvalues", [])):
                                cell = _unwrap_flowable(r[0])
                                p_line = (
                                    tf.paragraphs[0]
                                    if r_idx == 0
                                    else tf.add_paragraph()
                                )

                                left_indent = 0.0
                                if isinstance(cell, Paragraph):
                                    left_indent = getattr(cell.style, "leftIndent", 0.0)
                                    num_spaces = 0
                                    if left_indent > 0:
                                        num_spaces = max(
                                            0, int(round(left_indent / 4.8)) - 4
                                        )
                                    indent_str = " " * num_spaces

                                    if indent_str:
                                        run_ind = p_line.add_run()
                                        run_ind.text = indent_str
                                        run_ind.font.name = "Courier New"
                                        run_ind.font.size = Pt(10)
                                        run_ind.font.color.rgb = text_rgb

                                    parser = PptxHtmlParser(
                                        p_line,
                                        default_font_name="Courier New",
                                        default_font_size=10,
                                        default_color=text_rgb,
                                    )
                                    parser.feed(getattr(cell, "text", ""))
                                else:
                                    parser = PptxHtmlParser(
                                        p_line,
                                        default_font_name="Courier New",
                                        default_font_size=10,
                                        default_color=text_rgb,
                                    )
                                    parser.feed(str(cell))

                            y_cursor += h_tbl
                        elif is_callout:
                            text_val = _get_val_plain(
                                _unwrap_flowable(item._cellvalues[0][0])
                            )
                            num_lines = math.ceil(len(text_val) / 80.0)
                            h_tbl = num_lines * 0.3 + 0.4
                            tb = slide.shapes.add_textbox(
                                Inches(0.8),
                                Inches(y_cursor),
                                Inches(11.73),
                                Inches(h_tbl),
                            )
                            tf = tb.text_frame
                            tf.word_wrap = True
                            p_callout = tf.paragraphs[0]

                            parser = PptxHtmlParser(
                                p_callout,
                                default_font_name=t_theme.body_font,
                                default_font_size=12,
                                default_color=text_rgb,
                            )
                            cell = _unwrap_flowable(item._cellvalues[0][0])
                            if type(cell).__name__ == "Paragraph":
                                parser.feed(getattr(cell, "text", ""))
                            else:
                                parser.feed(str(cell))

                            fill = tb.fill
                            fill.solid()
                            if is_tip:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.green_bg)
                                tb.line.color.rgb = hex_to_rgb(t_theme.green)
                            elif is_note:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.yellow_bg)
                                tb.line.color.rgb = hex_to_rgb(t_theme.yellow)
                            elif is_warning:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.red_bg)
                                tb.line.color.rgb = hex_to_rgb(t_theme.red)
                            elif is_important:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.purple_bg)
                                tb.line.color.rgb = hex_to_rgb(t_theme.purple)
                            elif is_exam:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.yellow_bg)
                                tb.line.color.rgb = hex_to_rgb(t_theme.yellow)
                            elif is_theorem:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.purple_bg)
                                tb.line.color.rgb = hex_to_rgb(t_theme.purple)
                            elif is_definition:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.surface_alt)
                                tb.line.color.rgb = hex_to_rgb(t_theme.accent)
                            elif is_highlight:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.surface_alt)
                                tb.line.color.rgb = hex_to_rgb(t_theme.yellow)
                            elif is_question:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.surface_alt)
                                tb.line.color.rgb = hex_to_rgb(t_theme.accent)
                            elif is_answer:
                                fill.fore_color.rgb = hex_to_rgb(t_theme.surface)
                                tb.line.color.rgb = hex_to_rgb(t_theme.green)

                            tb.line.width = Pt(1.5)
                            y_cursor += h_tbl
                        elif is_flashcard:
                            inner_content = _unwrap_flowable(
                                getattr(item, "_cellvalues", [])[0][0]
                            )
                            inner_table = (
                                next(
                                    (
                                        x
                                        for x in inner_content
                                        if type(x).__name__ == "Table"
                                    ),
                                    None,
                                )
                                if isinstance(inner_content, list)
                                else (
                                    inner_content
                                    if type(inner_content).__name__ == "Table"
                                    else None
                                )
                            )
                            if inner_table is not None:
                                inner_cell = getattr(inner_table, "_cellvalues", [])[0][
                                    0
                                ]
                                text_val = _get_val_plain(_unwrap_flowable(inner_cell))
                            else:
                                text_val = ""
                            num_lines = math.ceil(len(text_val) / 80.0)
                            h_tbl = num_lines * 0.3 + 0.4
                            tb = slide.shapes.add_textbox(
                                Inches(0.8),
                                Inches(y_cursor),
                                Inches(11.73),
                                Inches(h_tbl),
                            )
                            tf = tb.text_frame
                            tf.word_wrap = True
                            p_fc = tf.paragraphs[0]
                            p_fc.text = text_val
                            p_fc.font.size = Pt(12)
                            p_fc.font.color.rgb = text_rgb
                            p_fc.font.name = t_theme.body_font

                            fill = tb.fill
                            fill.solid()
                            fill.fore_color.rgb = hex_to_rgb(t_theme.surface_alt)
                            tb.line.color.rgb = hex_to_rgb(t_theme.accent)
                            tb.line.width = Pt(2.0)
                            y_cursor += h_tbl
                        elif is_mcq:
                            mcq_lines = []
                            for r in getattr(item, "_cellvalues", []):
                                cell = _unwrap_flowable(r[0])
                                if isinstance(cell, Paragraph):
                                    mcq_lines.append(_get_val_plain(cell))
                            mcq_text = "\n".join(mcq_lines)
                            num_lines = len(mcq_lines) + 2
                            h_tbl = num_lines * 0.25 + 0.4
                            tb = slide.shapes.add_textbox(
                                Inches(0.8),
                                Inches(y_cursor),
                                Inches(11.73),
                                Inches(h_tbl),
                            )
                            tf = tb.text_frame
                            tf.word_wrap = True
                            p_mcq = tf.paragraphs[0]
                            p_mcq.text = mcq_text
                            p_mcq.font.size = Pt(11)
                            p_mcq.font.color.rgb = text_rgb
                            p_mcq.font.name = t_theme.body_font

                            fill = tb.fill
                            fill.solid()
                            fill.fore_color.rgb = hex_to_rgb(t_theme.surface_alt)
                            tb.line.color.rgb = hex_to_rgb(t_theme.accent2)
                            tb.line.width = Pt(1.5)
                            y_cursor += h_tbl
                        else:
                            is_packet = getattr(item, "_is_packet_format", False)
                            is_frame = getattr(item, "_is_frame_format", False)
                            rows_cnt = len(item._cellvalues)
                            cols_cnt = len(item._cellvalues[0]) if rows_cnt > 0 else 0
                            if rows_cnt > 0 and cols_cnt > 0:
                                h_tbl = rows_cnt * (0.3 if is_packet else 0.4) + 0.3
                                table_shape = slide.shapes.add_table(
                                    rows_cnt,
                                    cols_cnt,
                                    Inches(0.8),
                                    Inches(y_cursor),
                                    Inches(11.73),
                                    Inches(h_tbl),
                                )
                                tbl = table_shape.table

                                if (
                                    not hasattr(item, "_spanRanges")
                                    or not item._spanRanges
                                ):
                                    try:
                                        item._calcSpanRanges()
                                    except Exception:
                                        pass
                                spans_dict = getattr(item, "_spanRanges", {})

                                if spans_dict:
                                    for coord, span_val in spans_dict.items():
                                        if span_val is not None:
                                            sc, sr, ec, er = span_val
                                            if ec < cols_cnt and er < rows_cnt:
                                                cell_start = tbl.cell(sr, sc)
                                                cell_end = tbl.cell(er, ec)
                                                cell_start.merge(cell_end)

                                for r_idx in range(rows_cnt):
                                    for c_idx in range(cols_cnt):
                                        if (
                                            spans_dict
                                            and spans_dict.get((c_idx, r_idx)) is None
                                            and (c_idx, r_idx) in spans_dict
                                        ):
                                            continue

                                        cell = tbl.cell(r_idx, c_idx)
                                        val = _unwrap_flowable(
                                            item._cellvalues[r_idx][c_idx]
                                        )
                                        plain_cell = _get_val_plain(val)
                                        cell.text = plain_cell
                                        cell.fill.solid()

                                        if is_packet:
                                            if r_idx == 0 and len(item._cellvalues) > 1:
                                                cell.fill.fore_color.rgb = hex_to_rgb(
                                                    t_theme.surface_alt
                                                )
                                                font_sz = Pt(7)
                                            else:
                                                cell.fill.fore_color.rgb = surface_rgb
                                                font_sz = Pt(8)
                                        elif is_frame:
                                            cell.fill.fore_color.rgb = surface_rgb
                                            font_sz = Pt(9)
                                        else:
                                            if r_idx == 0:
                                                cell.fill.fore_color.rgb = hex_to_rgb(
                                                    t_theme.table_hdr
                                                )
                                                font_sz = Pt(11)
                                            else:
                                                bg_hex = (
                                                    t_theme.surface
                                                    if r_idx % 2 == 1
                                                    else t_theme.surface_alt
                                                )
                                                cell.fill.fore_color.rgb = hex_to_rgb(
                                                    bg_hex
                                                )
                                                font_sz = Pt(10)

                                        p_cell = cell.text_frame.paragraphs[0]
                                        p_cell.font.size = font_sz
                                        p_cell.font.color.rgb = text_rgb
                                        p_cell.font.name = t_theme.body_font

                                y_cursor += h_tbl

                    elif item_name == "Spacer":
                        y_cursor += getattr(item, "height", 0.0) / 72.0

                    elif (
                        item.__class__.__name__ == "Drawing"
                        or hasattr(item, "drawing")
                        or item_name == "ResponsiveDrawingFlowable"
                    ):
                        drawing = (
                            getattr(item, "drawing", None)
                            if hasattr(item, "drawing")
                            else item
                        )
                        if drawing is not None:
                            h_img = getattr(drawing, "height", 0.0) / 72.0
                            w_img = getattr(drawing, "width", 0.0) / 72.0

                            max_h = 4.5
                            max_w = 11.5
                            if h_img > max_h:
                                scale = max_h / h_img
                                h_img *= scale
                                w_img *= scale
                            if w_img > max_w:
                                scale = max_w / w_img
                                h_img *= scale
                                w_img *= scale

                            from reportlab.graphics import renderPM

                            temp_png = f"temp_slide_diag_{g_idx}_{y_cursor:.2f}.png"
                            success = False
                            try:
                                renderPM.drawToFile(drawing, temp_png, fmt="PNG")
                                success = True
                            except Exception:
                                temp_pdf = temp_png.replace(".png", ".pdf")
                                try:
                                    from reportlab.graphics import renderPDF
                                    import fitz

                                    renderPDF.drawToFile(drawing, temp_pdf)
                                    pdf_doc = fitz.open(temp_pdf)
                                    page = pdf_doc[0]
                                    pix = page.get_pixmap(dpi=300)  # type: ignore
                                    pix.save(temp_png)
                                    pdf_doc.close()
                                    success = True
                                except Exception:
                                    pass
                                finally:
                                    if os.path.exists(temp_pdf):
                                        try:
                                            os.remove(temp_pdf)
                                        except Exception:
                                            pass

                            if success and os.path.exists(temp_png):
                                try:
                                    x_pos = 0.8 + (11.73 - w_img) / 2.0
                                    slide.shapes.add_picture(
                                        temp_png,
                                        Inches(x_pos),
                                        Inches(y_cursor),
                                        width=Inches(w_img),
                                    )
                                    y_cursor += h_img + 0.2
                                except Exception:
                                    pass
                                finally:
                                    if os.path.exists(temp_png):
                                        try:
                                            os.remove(temp_png)
                                        except Exception:
                                            pass

                    elif item_name == "ImageFlowable":
                        local_path = getattr(item, "local_path", "")
                        caption = getattr(item, "caption", "")
                        if local_path and os.path.exists(local_path):
                            w_img = getattr(item, "w", 360.0) / 72.0
                            h_img = getattr(item, "h", 270.0) / 72.0

                            max_h = 4.5
                            max_w = 11.5
                            if h_img > max_h:
                                scale = max_h / h_img
                                h_img *= scale
                                w_img *= scale
                            if w_img > max_w:
                                scale = max_w / w_img
                                h_img *= scale
                                w_img *= scale

                            try:
                                x_pos = 0.8 + (11.73 - w_img) / 2.0
                                slide.shapes.add_picture(
                                    local_path,
                                    Inches(x_pos),
                                    Inches(y_cursor),
                                    width=Inches(w_img),
                                )
                                y_cursor += h_img + 0.2
                            except Exception:
                                pass
                        else:
                            # Render a styled missing image placeholder block in PowerPoint slides
                            w_img = getattr(item, "w", 360.0) / 72.0
                            h_img = getattr(item, "h", 270.0) / 72.0

                            max_h = 3.5
                            max_w = 10.0
                            if h_img > max_h:
                                scale = max_h / h_img
                                h_img *= scale
                                w_img *= scale
                            if w_img > max_w:
                                scale = max_w / w_img
                                h_img *= scale
                                w_img *= scale

                            try:
                                x_pos = 0.8 + (11.73 - w_img) / 2.0
                                shape = slide.shapes.add_shape(
                                    MSO_SHAPE.RECTANGLE,
                                    Inches(x_pos),
                                    Inches(y_cursor),
                                    Inches(w_img),
                                    Inches(h_img),
                                )
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = surface_rgb
                                shape.line.color.rgb = hex_to_rgb(t_theme.table_bdr)
                                shape.line.width = Pt(1)

                                tf = shape.text_frame
                                tf.word_wrap = True
                                p_tf = tf.paragraphs[0]
                                p_tf.text = f"[Image Not Available]\n{caption or ''}"
                                p_tf.font.size = Pt(12)
                                p_tf.font.color.rgb = text_rgb
                                p_tf.font.name = t_theme.body_font
                                p_tf.alignment = PP_ALIGN.CENTER

                                y_cursor += h_img + 0.2
                            except Exception:
                                pass

    prs.save(filename)


def build_html(
    output_dir,
    story=None,
    title=None,
):
    from .helpers import story as default_story, math_latex_registry
    from .theme import get_theme
    from reportlab.graphics import renderSVG
    import os
    import re

    if story is None:
        story = default_story

    t_theme = get_theme()

    def _fix_svg_scaling(svg_str):
        match_svg = re.search(r"<svg([^>]*)>", svg_str)
        if not match_svg:
            return svg_str
        attrs_str = match_svg.group(1)
        match_w = re.search(r'\bwidth\s*=\s*[\'"]([^\'"]+)[\'"]', attrs_str)
        match_h = re.search(r'\bheight\s*=\s*[\'"]([^\'"]+)[\'"]', attrs_str)
        if match_w and match_h:
            w_val = match_w.group(1).replace("pt", "").replace("px", "").strip()
            h_val = match_h.group(1).replace("pt", "").replace("px", "").strip()
            try:
                w = float(w_val)
                h = float(h_val)
                attrs_str = re.sub(r'\bwidth\s*=\s*[\'"][^\'"]*[\'"]', "", attrs_str)
                attrs_str = re.sub(r'\bheight\s*=\s*[\'"][^\'"]*[\'"]', "", attrs_str)
                attrs_str = re.sub(r'\bviewBox\s*=\s*[\'"][^\'"]*[\'"]', "", attrs_str)
                attrs_str = re.sub(r'\bstyle\s*=\s*[\'"][^\'"]*[\'"]', "", attrs_str)
                attrs_str = re.sub(
                    r'\bpreserveAspectRatio\s*=\s*[\'"][^\'"]*[\'"]', "", attrs_str
                )
                attrs_str = " ".join(attrs_str.split()).strip()

                new_tag = f'<svg {attrs_str} viewBox="0 0 {w} {h}" style="width: 100%; max-width: {w}px; height: auto; display: block; margin: 0 auto; overflow: visible;">'
                svg_str = (
                    svg_str[: match_svg.start()] + new_tag + svg_str[match_svg.end() :]
                )
            except ValueError:
                pass
        return svg_str

    def _get_val_html(val):
        val = _unwrap_flowable(val)
        if val is None:
            return ""
        if isinstance(val, str):
            return val
        cls_n = val.__class__.__name__
        if cls_n == "Spacer":
            return ""
        if cls_n == "Paragraph":
            return getattr(val, "text", "")
        if cls_n == "LaTeXFlowable":
            return f"$${getattr(val, 'latex_str', '')}$$"
        if cls_n == "ListFlowable":
            content_html = [
                '<ul style="padding-left: 20px; margin: 10px 0; text-align: left;">'
            ]
            items = getattr(
                val,
                "_flowables",
                getattr(val, "_content", getattr(val, "_list_content", [])),
            )
            for li in items:
                content_html.append(
                    f'<li style="margin-bottom: 6px;">{_get_val_html(li)}</li>'
                )
            content_html.append("</ul>")
            return "".join(content_html)
        if (
            cls_n == "Drawing"
            or hasattr(val, "drawing")
            or cls_n == "ResponsiveDrawingFlowable"
        ):
            drawing = getattr(val, "drawing", None) if hasattr(val, "drawing") else val
            if (
                drawing is not None
                and getattr(getattr(drawing, "__class__", None), "__name__", "")
                == "Drawing"
            ):
                try:
                    svg_str = renderSVG.drawToString(drawing)  # type: ignore
                    svg_clean = re.sub(r"<\?xml[^>]*\?>", "", svg_str)
                    svg_clean = re.sub(r"<!DOCTYPE[^>]*>", "", svg_clean)
                    svg_scaled = _fix_svg_scaling(svg_clean)
                    return f'<div style="text-align: center; width: 100%; cursor: zoom-in;" onclick="openDiagramOverlay(this)">{svg_scaled}</div>'
                except Exception:
                    pass
        if isinstance(val, (list, tuple)):
            return "".join(_get_val_html(x) for x in val)
        return str(val)

    os.makedirs(output_dir, exist_ok=True)

    sidebar_links = []
    content_html = []
    bookmark_counter = 0
    diagram_counter = 0

    flattened_story = []
    for item in story:
        if item.__class__.__name__ == "KeepTogether" and hasattr(item, "_content"):
            flattened_story.extend(getattr(item, "_content", []))
        else:
            flattened_story.append(item)

    in_theme_section = False

    def start_theme_section(theme):
        nonlocal in_theme_section
        res = []
        if in_theme_section:
            res.append("</div></div>")
        res.append(
            f'<div class="theme-section" style="background-color: {theme.bg}; color: {theme.text}; width: 100%; box-sizing: border-box; padding: 40px 0; min-height: 100vh; transition: background-color 0.3s, color 0.3s;">'
            f'<div style="max-width: 800px; margin: 0 auto; padding: 0 40px;">'
        )
        in_theme_section = True
        return "".join(res)

    for item in flattened_story:
        if getattr(item, "_is_toc_element", False):
            continue

        cls_name = item.__class__.__name__

        if cls_name == "ThemeSetterFlowable":
            t_theme = item.theme
            content_html.append(start_theme_section(t_theme))
            continue

        if not in_theme_section:
            content_html.append(start_theme_section(t_theme))

        if getattr(item, "_is_part_box", False):
            title_text = getattr(item, "_box_title", "")
            plain_title = re.sub(r"<[^>]+>", "", title_text)
            anchor = f"part-{bookmark_counter}"
            bookmark_counter += 1
            sidebar_links.append(
                f'<a href="#{anchor}" class="nav-part">{plain_title}</a>'
            )
            content_html.append(
                f'<div id="{anchor}" class="part-header" style="border-bottom: 2px solid {t_theme.table_bdr}; margin-top: 30px; padding-bottom: 10px;"><h1 style="color: {t_theme.accent}; margin: 0;">{title_text}</h1></div>'
            )

        elif getattr(item, "_is_chap_box", False):
            title_text = getattr(item, "_box_title", "")
            plain_title = re.sub(r"<[^>]+>", "", title_text)
            anchor = f"chap-{bookmark_counter}"
            bookmark_counter += 1
            sidebar_links.append(
                f'<a href="#{anchor}" class="nav-chap">{plain_title}</a>'
            )
            content_html.append(
                f'<div id="{anchor}" class="chap-header" style="border-bottom: 2px solid {t_theme.table_bdr}; margin-top: 30px; padding-bottom: 10px;"><h2 style="color: {t_theme.accent}; margin: 0;">{title_text}</h2></div>'
            )

        elif getattr(item, "_is_section", False):
            title_text = getattr(item, "_section_title", "")
            plain_title = re.sub(r"<[^>]+>", "", title_text)
            anchor = f"sect-{bookmark_counter}"
            bookmark_counter += 1
            sidebar_links.append(
                f'<a href="#{anchor}" class="nav-sect">{plain_title}</a>'
            )
            content_html.append(
                f'<h3 id="{anchor}" class="section-title" style="color: {t_theme.accent}; margin-top: 24px; margin-bottom: 12px;">{title_text}</h3><hr class="section-rule" style="border: 0; border-top: 1px solid {t_theme.accent}; margin-bottom: 15px;" />'
            )

        elif getattr(item, "_is_subsection", False):
            title_text = getattr(item, "_subsection_title", "")
            content_html.append(
                f'<h4 class="subsection-title" style="color: {t_theme.accent}; margin-top: 24px; margin-bottom: 12px;">{title_text}</h4>'
            )

        elif cls_name == "Paragraph":
            p_text = getattr(item, "text", "")
            content_html.append(f'<p class="body-paragraph">{p_text}</p>')

        elif cls_name == "IndexPrinterFlowable":
            content_html.append(f'<p class="body-paragraph" style="color: {t_theme.text_dim}; font-style: italic;">Index mapping is only available in the PDF format.</p>')

        elif cls_name in ("Table", "CodeBlockTable"):
            is_code = getattr(item, "_is_code_block", False)
            is_tip = getattr(item, "_is_tip", False)
            is_note = getattr(item, "_is_note", False)
            is_warning = getattr(item, "_is_warning", False)
            is_important = getattr(item, "_is_important", False)
            is_exam = getattr(item, "_is_exam", False)
            is_theorem = getattr(item, "_is_theorem", False)
            is_definition = getattr(item, "_is_definition", False)
            is_highlight = getattr(item, "_is_highlight", False)
            is_question = getattr(item, "_is_question", False)
            is_answer = getattr(item, "_is_answer", False)
            is_mcq = getattr(item, "_is_mcq", False)
            is_flashcard = getattr(item, "_is_flashcard", False)
            is_frame_format = getattr(item, "_is_frame_format", False)
            is_packet_format = getattr(item, "_is_packet_format", False)

            if is_code and len(item._cellvalues) > 0:
                code_lines = []
                for r in item._cellvalues:
                    cell = _unwrap_flowable(r[0])
                    if isinstance(cell, Paragraph):
                        left_indent = getattr(cell.style, "leftIndent", 0.0)
                        num_spaces = 0
                        if left_indent > 0:
                            num_spaces = max(0, int(round(left_indent / 4.8)) - 4)
                        indent_str = " " * num_spaces
                        line_text = indent_str + getattr(cell, "text", "")
                        code_lines.append(line_text)
                    else:
                        code_lines.append(str(cell))
                code_text = "\n".join(code_lines)
                content_html.append(
                    f'<div class="code-container" style="position: relative; margin: 15px 0;">'
                    f'<button class="copy-btn" onclick="copyCode(this)" style="position: absolute; right: 10px; top: 10px; background: {t_theme.surface_alt}; border: 1px solid {t_theme.table_bdr}; color: {t_theme.text}; padding: 4px 8px; border-radius: 4px; font-size: 0.75em; cursor: pointer; opacity: 0; transition: opacity 0.2s; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">Copy</button>'
                    f'<pre style="background-color: {t_theme.code_bg}; color: {t_theme.text_code}; padding: 16px; border-radius: 4px; border-left: 4px solid {t_theme.accent}; overflow-x: auto; font-family: Courier, monospace; margin: 0; white-space: pre;"><code>{code_text}</code></pre>'
                    f"</div>"
                )
            elif is_tip:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout tip" style="border-left: 4px solid {t_theme.green}; background-color: {t_theme.green_bg}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_note:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout note" style="border-left: 4px solid {t_theme.yellow}; background-color: {t_theme.yellow_bg}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_warning:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout warning" style="border-left: 4px solid {t_theme.red}; background-color: {t_theme.red_bg}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_important:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout important" style="border-left: 4px solid {t_theme.purple}; background-color: {t_theme.purple_bg}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_exam:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout exam" style="border-left: 4px solid {t_theme.yellow}; background-color: {t_theme.yellow_bg}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_theorem:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout theorem" style="border: 1px solid {t_theme.purple}; background-color: {t_theme.purple_bg}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_definition:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout definition" style="border-left: 4px solid {t_theme.accent}; background-color: {t_theme.surface_alt}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_highlight:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout highlight" style="border-left: 4px solid {t_theme.yellow}; background-color: {t_theme.surface_alt}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_question:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout question" style="border-left: 4px solid {t_theme.accent}; background-color: {t_theme.surface_alt}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_answer:
                text = _get_val_html(_unwrap_flowable(item._cellvalues[0][0]))
                content_html.append(
                    f'<div class="callout answer" style="border-left: 4px solid {t_theme.green}; background-color: {t_theme.surface}; color: {t_theme.text}; padding: 12px; margin: 15px 0; border-radius: 4px;">{text}</div>'
                )
            elif is_flashcard:
                inner_content = _unwrap_flowable(getattr(item, "_cellvalues", [])[0][0])
                inner_table = (
                    next(
                        (x for x in inner_content if type(x).__name__ == "Table"),
                        None,
                    )
                    if isinstance(inner_content, list)
                    else (
                        inner_content
                        if type(inner_content).__name__ == "Table"
                        else None
                    )
                )
                if inner_table is not None:
                    inner_cell = getattr(inner_table, "_cellvalues", [])[0][0]
                    text = _get_val_html(_unwrap_flowable(inner_cell))
                else:
                    text = ""
                content_html.append(
                    f'<div class="callout flashcard" style="border: 2px solid {t_theme.accent}; background-color: {t_theme.surface_alt}; color: {t_theme.text}; padding: 16px; margin: 15px 0; border-radius: 8px; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">{text}</div>'
                )
            elif is_mcq:
                mcq_html = [
                    f'<div class="callout mcq" style="border-left: 4px solid {t_theme.accent2}; background-color: {t_theme.surface_alt}; color: {t_theme.text}; padding: 16px; margin: 15px 0; border-radius: 4px;">'
                ]
                for r in getattr(item, "_cellvalues", []):
                    cell = _unwrap_flowable(r[0])
                    if isinstance(cell, Paragraph):
                        text = _get_val_html(cell)
                        mcq_html.append(f'<p style="margin: 6px 0;">{text}</p>')
                mcq_html.append("</div>")
                content_html.append("".join(mcq_html))
            else:
                is_packet = is_packet_format
                is_frame = is_frame_format

                if is_packet:
                    tbl_style = f"width: 100%; max-width: 600px; table-layout: fixed; border-collapse: collapse; margin: 15px auto; border: 1px solid {t_theme.table_bdr}; font-size: 0.65em;"
                elif is_frame:
                    tbl_style = f"width: 100%; max-width: 600px; table-layout: fixed; border-collapse: collapse; margin: 15px auto; border: 1px solid {t_theme.table_bdr}; font-size: 0.8em;"
                else:
                    tbl_style = f"width: 100%; max-width: 650px; border-collapse: collapse; margin: 15px auto; border: 1px solid {t_theme.table_bdr};"

                if not hasattr(item, "_spanRanges") or not item._spanRanges:
                    try:
                        item._calcSpanRanges()
                    except Exception:
                        pass
                spans_dict = getattr(item, "_spanRanges", {})

                html_tbl = [
                    f'<div class="table-container" style="overflow-x: auto; width: 100%; text-align: center;">'
                    f'<table style="{tbl_style}">'
                ]
                for r_idx, r in enumerate(getattr(item, "_cellvalues", [])):
                    bg = (
                        t_theme.table_hdr
                        if r_idx == 0 and not is_packet and not is_frame
                        else (
                            t_theme.surface if r_idx % 2 == 1 else t_theme.surface_alt
                        )
                    )
                    txt_color = (
                        "#ffffff"
                        if r_idx == 0 and not is_packet and not is_frame
                        else t_theme.text
                    )

                    if is_packet:
                        if r_idx == 0 and len(getattr(item, "_cellvalues", [])) > 1:
                            bg = t_theme.surface_alt
                            txt_color = t_theme.text_dim

                    html_tbl.append(
                        f'<tr style="background-color: {bg}; color: {txt_color}; border-bottom: 1px solid {t_theme.table_bdr};">'
                    )
                    for c_idx, val in enumerate(r):
                        if (
                            spans_dict
                            and (c_idx, r_idx) in spans_dict
                            and spans_dict[(c_idx, r_idx)] is None
                        ):
                            continue
                        tag = (
                            "th"
                            if r_idx == 0 and not is_packet and not is_frame
                            else "td"
                        )

                        colspan_attr = ""
                        rowspan_attr = ""
                        if spans_dict and (c_idx, r_idx) in spans_dict:
                            span_val = spans_dict[(c_idx, r_idx)]
                            if span_val:
                                sc, sr, ec, er = span_val
                                cs = ec - sc + 1
                                rs = er - sr + 1
                                if cs > 1:
                                    colspan_attr = f' colspan="{cs}"'
                                if rs > 1:
                                    rowspan_attr = f' rowspan="{rs}"'

                        if is_packet:
                            cell_style = f"padding: 4px 1px; text-align: center; border: 1px solid {t_theme.table_bdr}; overflow: hidden;"
                        elif is_frame:
                            cell_style = f"padding: 6px 4px; text-align: center; border: 1px solid {t_theme.table_bdr};"
                        else:
                            cell_style = f"padding: 8px 12px; text-align: left; border: 1px solid {t_theme.table_bdr};"

                        html_tbl.append(
                            f'<{tag}{colspan_attr}{rowspan_attr} style="{cell_style}">{_get_val_html(val)}</{tag}>'
                        )
                    html_tbl.append("</tr>")
                html_tbl.append("</table></div>")
                content_html.append("".join(html_tbl))

        elif (
            item.__class__.__name__ == "Drawing"
            or hasattr(item, "drawing")
            or cls_name == "ResponsiveDrawingFlowable"
        ):
            drawing = (
                getattr(item, "drawing", None) if hasattr(item, "drawing") else item
            )
            if (
                drawing is not None
                and getattr(getattr(drawing, "__class__", None), "__name__", "")
                == "Drawing"
            ):
                try:
                    svg_str = renderSVG.drawToString(drawing)  # type: ignore
                    svg_clean = re.sub(r"<\?xml[^>]*\?>", "", svg_str)
                    svg_clean = re.sub(r"<!DOCTYPE[^>]*>", "", svg_clean)

                    # Make clip path and group IDs unique to avoid ID collision
                    clip_id = f"clip_{diagram_counter}"
                    group_id = f"group_{diagram_counter}"
                    svg_clean = svg_clean.replace('id="clip"', f'id="{clip_id}"')
                    svg_clean = svg_clean.replace("url(#clip)", f"url(#{clip_id})")
                    svg_clean = svg_clean.replace('id="group"', f'id="{group_id}"')
                    diagram_counter += 1

                    svg_scaled = _fix_svg_scaling(svg_clean)

                    # Retrieve the diagram theme if attached to drawing or item
                    diag_theme = getattr(
                        item, "diagram_theme", getattr(drawing, "_diagram_theme", None)
                    )

                    # Check dark/light mode and apply color inversion filter if they differ
                    def is_dark_hex(hex_color):
                        h = hex_color.lstrip("#")
                        if len(h) == 3:
                            h = "".join(c * 2 for c in h)
                        try:
                            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
                            return luminance < 0.5
                        except Exception:
                            return True

                    page_is_dark = is_dark_hex(t_theme.bg)
                    diag_is_dark = (
                        is_dark_hex(diag_theme.bg) if diag_theme else page_is_dark
                    )

                    filter_style = ""
                    if page_is_dark != diag_is_dark:
                        filter_style = "filter: invert(1) hue-rotate(180deg);"

                    # Always match the container background and border with the global theme to follow the whole page theme
                    bg_color = t_theme.surface
                    border_color = t_theme.table_bdr

                    content_html.append(
                        f'<div class="diagram-container" style="text-align: center; margin: 20px 0; background: {bg_color}; padding: 15px; border-radius: 6px; border: 1px solid {border_color}; overflow-x: auto; overflow-y: visible; cursor: zoom-in;" onclick="openDiagramOverlay(this)">'
                        f'<div style="{filter_style}">{svg_scaled}</div>'
                        f"</div>"
                    )
                except Exception:
                    pass

        elif cls_name == "ImageFlowable":
            src = getattr(item, "resolved_src", getattr(item, "src", ""))
            local_path = getattr(item, "local_path", "")
            caption = getattr(item, "caption", "")
            link = getattr(item, "link", None)

            # Check if image resolved/downloaded successfully on disk
            is_valid_image = local_path and os.path.exists(local_path)

            if is_valid_image:
                img_html = f'<img src="{src}" style="max-width: 100%; height: auto; display: block; margin: 0 auto; border-radius: 4px;" />'
                if link:
                    img_html = f'<a href="{link}" target="_blank">{img_html}</a>'
            else:
                w_fallback = getattr(item, "w", 150)
                h_fallback = getattr(item, "h", 100)
                
                # Limit length of display URL/path
                display_src = src
                if len(display_src) > 50:
                    display_src = display_src[:47] + "..."
                
                img_html = (
                    f'<div class="image-fallback-box" style="margin: 0 auto; max-width: {w_fallback}px; height: {h_fallback}px; '
                    f'border: 1px dashed {t_theme.table_bdr}; background: {t_theme.surface}; border-radius: 4px; '
                    f'display: flex; flex-direction: column; align-items: center; justify-content: center; '
                    f'color: {t_theme.text_dim}; font-size: 0.9em; padding: 10px; box-sizing: border-box; text-align: center; line-height: 1.2;">'
                    f'<span style="font-size: 1.4em; margin-bottom: 4px;">⚠️</span>'
                    f'<strong>[Image Not Available]</strong>'
                    f'<span style="font-size: 0.75em; margin-top: 4px; word-break: break-all; opacity: 0.85;">{display_src}</span>'
                    f'</div>'
                )
                if link:
                    img_html = f'<a href="{link}" target="_blank" style="text-decoration: none; display: inline-block;">{img_html}</a>'

            if caption:
                img_html += f'<p style="text-align: center; font-style: italic; font-size: 0.9em; margin-top: 5px; color: {t_theme.text_dim};">{caption}</p>'

            content_html.append(
                f'<div class="image-container" style="text-align: center; margin: 20px 0;">'
                f'{img_html}'
                f'</div>'
            )

        elif cls_name == "LaTeXFlowable":
            math_str = item.latex_str
            content_html.append(
                f'<div class="math-block" style="text-align: center; margin: 15px 0; font-size: 1.2em;">$${math_str}$$</div>'
            )

    if in_theme_section:
        content_html.append("</div></div>")

    extracted_title = None
    if not title:
        for item in flattened_story:
            if type(item).__name__ == "Table" and hasattr(item, "_cellvalues") and item._cellvalues:
                try:
                    cell = _unwrap_flowable(item._cellvalues[0][0])
                    if isinstance(cell, list):
                        cell = cell[0]
                    if type(cell).__name__ == "Paragraph":
                        p_style = getattr(cell, "style", None)
                        if p_style:
                            parent_name = getattr(getattr(p_style, "parent", None), "name", "NO_PARENT")
                            if parent_name in ("Cover_H1", "Chap_H1", "H1"):
                                raw_text = getattr(cell, "text", "")
                                clean_title = re.sub(r"<[^>]*>", "", raw_text).strip()
                                if clean_title:
                                    extracted_title = clean_title
                                    break
                except Exception:
                    pass

    doc_title = title or extracted_title or "Documentation"
    html_output = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{doc_title}</title>
    <script>
        window.MathJax = {{
            tex: {{
                inlineMath: [['$', '$'], ['\\\\(', '\\\\)']],
                displayMath: [['$$', '$$'], ['\\\\[', '\\\\]']]
            }},
            svg: {{
                fontCache: 'global'
            }}
        }};
    </script>
    <script id="MathJax-script" async src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-mml-chtml.js"></script>
    <script>
        function copyCode(button) {{
            const codeEl = button.nextElementSibling.querySelector('code');
            navigator.clipboard.writeText(codeEl.textContent).then(() => {{
                button.innerText = 'Copied!';
                setTimeout(() => {{
                    button.innerText = 'Copy';
                }}, 2000);
            }});
        }}
        function openDiagramOverlay(element) {{
            var overlay = document.getElementById('diagram-overlay');
            var content = document.getElementById('diagram-overlay-content');
            var svgHTML = element.innerHTML;
            content.innerHTML = svgHTML;
            overlay.style.display = 'flex';
        }}
        function toggleSidebar() {{
            const body = document.body;
            const btn = document.getElementById('sidebar-toggle');
            const icon = document.getElementById('toggle-icon');
            body.classList.toggle('sidebar-collapsed');
            
            if (body.classList.contains('sidebar-collapsed')) {{
                icon.innerText = '▶';
                btn.style.left = '15px';
            }} else {{
                icon.innerText = '◀';
                btn.style.left = '270px';
            }}
        }}
        window.addEventListener('DOMContentLoaded', () => {{
            if (window.innerWidth <= 768) {{
                toggleSidebar();
            }}
        }});
    </script>
    <style>
        :root {{
            --bg-color: {t_theme.bg};
            --text-color: {t_theme.text};
            --accent-color: {t_theme.accent};
            --sidebar-bg: {t_theme.surface};
            --border-color: {t_theme.table_bdr};
        }}
        body {{
            margin: 0;
            padding: 0;
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif;
            background-color: var(--bg-color);
            color: var(--text-color);
            display: flex;
            min-height: 100vh;
        }}
        html {{
            scrollbar-width: thin;
            scrollbar-color: var(--border-color) var(--bg-color);
        }}
        pre {{
            scrollbar-width: thin;
            scrollbar-color: var(--border-color) var(--bg-color);
        }}
        ::-webkit-scrollbar {{
            width: 8px;
            height: 8px;
        }}
        ::-webkit-scrollbar-track {{
            background: var(--bg-color);
        }}
        ::-webkit-scrollbar-thumb {{
            background: var(--border-color);
            border-radius: 4px;
        }}
        ::-webkit-scrollbar-thumb:hover {{
            background: var(--accent-color);
        }}
        pre::-webkit-scrollbar-track {{
            background: {t_theme.code_bg};
        }}
        pre::-webkit-scrollbar-thumb {{
            background: {t_theme.table_bdr};
        }}
        pre::-webkit-scrollbar-thumb:hover {{
            background: {t_theme.accent};
        }}
        .code-container:hover .copy-btn {{
            opacity: 1 !important;
        }}
        .copy-btn:hover {{
            background-color: var(--border-color) !important;
        }}
        @media (max-width: 768px) {{
            .copy-btn {{
                opacity: 1 !important;
            }}
        }}
        aside {{
            width: 260px;
            background-color: var(--sidebar-bg);
            border-right: 1px solid var(--border-color);
            padding: 20px 10px;
            box-sizing: border-box;
            position: fixed;
            top: 0;
            bottom: 0;
            overflow-y: auto;
            transition: transform 0.3s ease;
            z-index: 900;
        }}
        aside a {{
            display: block;
            color: var(--text-color);
            text-decoration: none;
            padding: 8px 12px;
            border-radius: 4px;
            margin-bottom: 4px;
            font-size: 0.9em;
            opacity: 0.8;
            transition: all 0.2s;
        }}
        aside a:hover {{
            background-color: var(--border-color);
            opacity: 1.0;
        }}
        aside a.nav-part {{
            font-weight: bold;
            font-size: 1.0em;
            margin-top: 15px;
            color: var(--accent-color);
        }}
        aside a.nav-chap {{
            font-weight: 600;
            padding-left: 20px;
        }}
        aside a.nav-sect {{
            padding-left: 32px;
            font-size: 0.85em;
            opacity: 0.75;
        }}
        main {{
            margin-left: 260px;
            padding: 0;
            width: calc(100% - 260px);
            min-height: 100vh;
            box-sizing: border-box;
            transition: margin-left 0.3s ease, width 0.3s ease;
        }}
        #sidebar-toggle {{
            position: fixed;
            left: 270px;
            top: 15px;
            z-index: 1000;
            background: var(--sidebar-bg);
            border: 1px solid var(--border-color);
            color: var(--text-color);
            padding: 8px 12px;
            border-radius: 4px;
            cursor: pointer;
            transition: left 0.3s ease, background-color 0.2s, color 0.2s;
            box-shadow: 0 2px 8px rgba(0,0,0,0.15);
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: bold;
        }}
        #sidebar-toggle:hover {{
            background: var(--border-color);
            color: var(--accent-color);
        }}
        body.sidebar-collapsed aside {{
            transform: translateX(-260px);
        }}
        body.sidebar-collapsed main {{
            margin-left: 0;
            width: 100%;
        }}
        h1, h2, h3, h4 {{
            color: var(--accent-color);
            margin-top: 24px;
            margin-bottom: 12px;
        }}
        .part-header, .chap-header {{
            border-bottom: 2px solid var(--border-color);
            padding-bottom: 10px;
            margin-top: 30px;
        }}
        .body-paragraph {{
            line-height: 1.6;
            margin-bottom: 16px;
        }}
        .diagram-container svg {{
            max-width: 100%;
            height: auto;
        }}
        @media (max-width: 768px) {{
            body {{
                flex-direction: column;
            }}
            aside {{
                width: 260px;
                position: fixed;
                border-right: 1px solid var(--border-color);
                border-bottom: none;
                max-height: none;
            }}
            main {{
                margin-left: 0 !important;
                width: 100% !important;
                padding: 0;
            }}
        }}
    </style>
</head>
<body>
    <button id="sidebar-toggle" onclick="toggleSidebar()">
        <span id="toggle-icon">◀</span>
    </button>
    <aside>
        <div style="font-weight: bold; font-size: 1.1em; padding: 10px 12px; color: var(--accent-color); border-bottom: 1px solid var(--border-color); margin-bottom: 15px;">
            {doc_title}
        </div>
        {"".join(sidebar_links)}
    </aside>
    <main>
        {"".join(content_html)}
    </main>
    <div id="diagram-overlay" style="display: none; position: fixed; top: 0; left: 0; width: 100vw; height: 100vh; background: rgba(0,0,0,0.85); z-index: 9999; align-items: center; justify-content: center; overflow: hidden; cursor: zoom-out;" onclick="this.style.display='none'">
        <div id="diagram-overlay-content" style="background: var(--bg-color); padding: 20px; border-radius: 8px; border: 1px solid var(--border-color); max-width: 95vw; max-height: 95vh; overflow: auto; display: flex; justify-content: center; align-items: center; cursor: default;" onclick="event.stopPropagation()">
        </div>
    </div>
</body>
</html>
"""

    # Convert math-inline images back to MathJax LaTeX strings
    def repl_math_inline(match):
        img_tag = match.group(0)
        src_match = re.search(r'src=["\']([^"\']+)["\']', img_tag)
        if src_match:
            src_path = src_match.group(1)
            # Match directly
            if src_path in math_latex_registry:
                latex = math_latex_registry[src_path]
                return f"${latex}$"
            # Match by base filename (resilient to path formatting)
            filename_part = os.path.basename(src_path)
            for k, v in math_latex_registry.items():
                if os.path.basename(k) == filename_part:
                    return f"${v}$"
        return img_tag

    html_output = re.sub(
        r"<img\b[^>]*math_[^>]*\.png[^>]*>", repl_math_inline, html_output
    )

    base_name = os.path.basename(os.path.normpath(output_dir))
    if base_name.endswith("_html"):
        base_name = base_name[:-5]
    elif base_name.endswith(".html"):
        base_name = base_name[:-5]
    
    html_filename = "index.html"
    
    with open(os.path.join(output_dir, html_filename), "w", encoding="utf-8") as f:
        f.write(html_output)


__all__ = [
    "PAGE_W",
    "PAGE_H",
    "PM",
    "CW",
    "page_decor",
    "build_doc",
    "ThemedCanvas",
    "LaTeXFlowable",
    "build_pptx",
    "build_html",
]
